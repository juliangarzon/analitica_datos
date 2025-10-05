#!/usr/bin/env python3
"""
Create PowerPoint presentation for Child Abuse Analysis - Palmira 2022
Based on class recommendations for effective data visualization
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("Error: python-pptx not available. Creating text-based presentation structure instead.")
    import sys
    sys.exit(1)

def create_presentation():
    # Create presentation
    prs = Presentation()

    # Define colors based on class recommendations
    colors = {
        'primary': RGBColor(31, 119, 180),     # Blue
        'danger': RGBColor(214, 39, 40),       # Red
        'warning': RGBColor(255, 127, 14),     # Orange
        'success': RGBColor(44, 160, 44),      # Green
        'dark': RGBColor(47, 79, 79),          # Dark gray
    }

    # SLIDE 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = 'CHILD ABUSE IN PALMIRA 2022'
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text = 'Data-Driven Analysis for Evidence-Based Policy\\n\\n📊 384 Official Cases Analyzed\\n📈 Validated Metrics & Insights\\n🎯 Actionable Recommendations'
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # SLIDE 2: Research Question & Premise
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'RESEARCH PREMISE TO VALIDATE'

    content = slide.placeholders[1]
    content.text = '''🎯 HYPOTHESIS:

"Child abuse in Palmira 2022 is characterized by being:

• Predominantly INTRAFAMILIAL (78% of aggressors are parents)
• MULTIDIMENSIONAL (91% involve multiple abuse types)
• Causes UNIVERSAL psychological damage (96% of cases)"


📋 METHODOLOGY:
• Statistical analysis of 384 official cases
• Validated metrics calculation
• Evidence-based conclusions'''

    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # SLIDE 3: Finding 1 - Intrafamilial Problem
    slide_layout = prs.slide_layouts[5]  # Title only
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'FINDING 1: INTRAFAMILIAL PROBLEM'

    # Add text box for key finding
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = '''✅ VALIDATED: 78.1% of aggressors are parents

📊 KEY METRICS:
• 300 out of 384 cases
• Mothers: 207 cases (53.9%)
• Fathers: 93 cases (24.2%)

🏠 IMPLICATION:
Family-centered intervention required'''

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)

    # Add placeholder for chart
    left = Inches(6)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(4.5)
    chart_placeholder = slide.shapes.add_textbox(left, top, width, height)
    chart_tf = chart_placeholder.text_frame
    chart_tf.text = '''[CHART PLACEHOLDER]

📊 Recommended visualization:
• Donut chart: Parents vs Others
• Bar chart: Madre, Padre, Other relatives

Color scheme:
• Red for parents (high risk)
• Orange for other family

Apply class principles:
• Clear hierarchy
• Minimal chartjunk
• Accessible colors'''

    for paragraph in chart_tf.paragraphs:
        paragraph.font.size = Pt(12)

    # SLIDE 4: Finding 2 - Multidimensional Abuse
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'FINDING 2: MULTIDIMENSIONAL ABUSE'

    # Add text box for key finding
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = '''✅ VALIDATED: 91.4% suffer multiple abuse types

📊 KEY METRICS:
• 351 out of 384 cases
• Average: 2.73 types per child
• Most common: 3 types simultaneously (245 cases)

🔄 IMPLICATION:
Holistic treatment approach needed'''

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)

    # Add placeholder for chart
    left = Inches(6)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(4.5)
    chart_placeholder = slide.shapes.add_textbox(left, top, width, height)
    chart_tf = chart_placeholder.text_frame
    chart_tf.text = '''[CHART PLACEHOLDER]

📊 Recommended visualization:
• Pie chart: Single vs Multiple types
• Bar chart: Distribution by number of types

Color scheme:
• Green for single (low complexity)
• Orange for 2 types
• Red for 3+ types (high complexity)

Apply class principles:
• Data-ink ratio maximization
• Meaningful color coding'''

    for paragraph in chart_tf.paragraphs:
        paragraph.font.size = Pt(12)

    # SLIDE 5: Finding 3 - Psychological Universality
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'FINDING 3: UNIVERSAL PSYCHOLOGICAL DAMAGE'

    # Add text box for key finding
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = '''✅ VALIDATED: 95.6% suffer psychological abuse

📊 ABUSE TYPE PREVALENCE:
• Psychological: 367 cases (95.6%)
• Verbal: 328 cases (85.4%)
• Physical: 277 cases (72.1%)
• Sexual: 47 cases (12.2%)
• Economic: 13 cases (3.4%)

🧠 IMPLICATION:
Mental health must be priority'''

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)

    # Add placeholder for chart
    left = Inches(6)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(4.5)
    chart_placeholder = slide.shapes.add_textbox(left, top, width, height)
    chart_tf = chart_placeholder.text_frame
    chart_tf.text = '''[CHART PLACEHOLDER]

📊 Recommended visualization:
• Horizontal bar chart by prevalence
• Highlight psychological abuse (biggest bar)

Color scheme:
• Red for high prevalence (80%+)
• Orange for medium (50-80%)
• Blue for lower prevalence

Apply class principles:
• Sort by magnitude
• Clear value labels'''

    for paragraph in chart_tf.paragraphs:
        paragraph.font.size = Pt(12)

    # SLIDE 6: Victim Demographics
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'VICTIM DEMOGRAPHICS: WHO IS AFFECTED?'

    # Add text box for demographics
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5)
    height = Inches(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = '''👥 GENDER BALANCE:
• Female: 203 cases (52.9%)
• Male: 181 cases (47.1%)

🎂 AGE DISTRIBUTION:
• Average age: 9.5 years
• Range: 0-18 years
• All age groups affected

📍 GEOGRAPHIC SPREAD:
• Multiple comunas affected
• Urban concentration pattern'''

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(16)

    # Add placeholder for chart
    left = Inches(6)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(4.5)
    chart_placeholder = slide.shapes.add_textbox(left, top, width, height)
    chart_tf = chart_placeholder.text_frame
    chart_tf.text = '''[CHART PLACEHOLDER]

📊 Recommended visualization:
• Pie chart: Gender distribution
• Histogram: Age distribution
• Add mean/median lines

Color scheme:
• Blue for male, Red for female
• Purple gradient for age histogram

Apply class principles:
• Multiple small multiples
• Clear axis labels'''

    for paragraph in chart_tf.paragraphs:
        paragraph.font.size = Pt(12)

    # SLIDE 7: Executive Dashboard
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'EXECUTIVE DASHBOARD: ALL FINDINGS'

    # Add instructions for dashboard
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(10)
    height = Inches(5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = '''[EXECUTIVE DASHBOARD PLACEHOLDER]

📊 Insert the comprehensive dashboard from the Jupyter notebook here

🎯 KEY ELEMENTS TO INCLUDE:
• 4 validated metrics as large numbers/gauges
• Abuse types prevalence (bar chart)
• Relationship breakdown (donut chart)
• Complexity distribution (bar chart)
• Key statistics summary box

🎨 DESIGN PRINCIPLES TO APPLY:
• Consistent color scheme (red for high severity, orange for medium, blue for neutral)
• Clear hierarchy (biggest numbers for most important metrics)
• Minimal text, maximum visual impact
• Professional layout suitable for executive presentation

📏 LAYOUT RECOMMENDATION:
• 2x3 grid layout
• Top row: 3 key metrics as large numbers
• Bottom row: 3 supporting visualizations
• Use white space effectively'''

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(14)

    # SLIDE 8: Policy Implications
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'POLICY IMPLICATIONS & RECOMMENDATIONS'

    content = slide.placeholders[1]
    content.text = '''🎯 EVIDENCE-BASED RECOMMENDATIONS:

🏠 FAMILY-CENTERED INTERVENTIONS (78% intrafamilial)
• Mandatory parent education programs
• Family counseling and support services
• Home-based intervention teams

🔄 INTEGRATED SERVICE DELIVERY (91% multiple abuse types)
• Holistic case management approach
• Cross-sector coordination (health, education, social services)
• Trauma-informed care protocols

🧠 MENTAL HEALTH PRIORITY (96% psychological damage)
• Specialized child psychologists
• Long-term therapy programs
• School-based mental health services

📊 INTENSIVE SUPPORT SYSTEM (2.7 avg abuse types)
• High-frequency monitoring
• Multi-disciplinary teams
• Crisis intervention protocols'''

    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)

    # SLIDE 9: Next Steps
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = 'NEXT STEPS: FROM DATA TO ACTION'

    content = slide.placeholders[1]
    content.text = '''🚀 IMMEDIATE ACTIONS:

📋 PROGRAM DEVELOPMENT:
• Design family intervention protocols
• Establish integrated service teams
• Create mental health referral pathways

📊 MONITORING & EVALUATION:
• Implement outcome tracking system
• Regular data collection and analysis
• Quarterly progress reviews

💰 RESOURCE ALLOCATION:
• Budget for specialized personnel
• Fund training programs
• Invest in technology platforms

🤝 STAKEHOLDER ENGAGEMENT:
• Present findings to city council
• Engage community organizations
• Collaborate with NGOs and universities

📅 TIMELINE: Implementation within 6 months'''

    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)

    # SLIDE 10: Thank You
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = 'THANK YOU'
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True

    subtitle.text = 'Questions & Discussion\\n\\n📊 Data-Driven Insights for Child Protection\\n🎯 Evidence-Based Policy Making\\n🚀 From Analysis to Action'
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save('maltrato_infantil_presentacion_base.pptx')
    print('✅ Presentation created successfully!')
    print('📁 File: maltrato_infantil_presentacion_base.pptx')
    print('📊 Ready for graphics insertion based on class principles')

if __name__ == "__main__":
    create_presentation()